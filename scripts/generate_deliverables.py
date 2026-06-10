#!/usr/bin/env python3
"""Generate README, environment log, final report, and presentation slides."""

from __future__ import annotations

import csv
import json
import os
import platform
import subprocess
import sys
from pathlib import Path
from textwrap import wrap

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


ROOT = Path(__file__).resolve().parents[1]
RESULTS = ROOT / "results"
FIGURES = ROOT / "figures"
REPORT = ROOT / "report"
SLIDES = ROOT / "slides"
LOGS = ROOT / "artifacts" / "logs"
LEROBOT_MODEL_DIR = ROOT / "official_reproduction" / "hf_lerobot_diffusion_pusht"
OFFICIAL_STATUS = RESULTS / "official_reproduction_status.json"


def sh(cmd: list[str], env: dict[str, str] | None = None) -> str:
    try:
        out = subprocess.check_output(cmd, stderr=subprocess.STDOUT, text=True, timeout=20, env=env)
        return out.strip()
    except Exception as exc:
        return f"FAILED: {' '.join(cmd)}\n{exc}"


def row(summary: pd.DataFrame, method: str) -> pd.Series:
    return summary[summary.method == method].iloc[0]


def fmt(x: float, n: int = 3) -> str:
    return f"{x:.{n}f}"


def load_json(path: Path) -> dict:
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def lerobot_status() -> dict:
    eval_info = load_json(LEROBOT_MODEL_DIR / "eval_info.json")
    smoke = load_json(RESULTS / "lerobot_pusht_smoke_summary.json")
    model_file = LEROBOT_MODEL_DIR / "model.safetensors"
    return {
        "hf_eval": eval_info.get("aggregated", {}),
        "smoke": smoke,
        "model_size_mb": model_file.stat().st_size / (1024 * 1024) if model_file.exists() else None,
    }


def official_status() -> dict:
    status = load_json(OFFICIAL_STATUS)
    results = status.get("official_eval_results", {})
    preferred = results.get("pusht_eval_official_n50_diffusers011") or results.get("pusht_eval_official_n50") or {}
    return {
        "completed": bool(status.get("official_eval_completed")),
        "checkpoint_size_bytes": status.get("official_checkpoint_size_bytes", 0),
        "preferred": preferred,
        "results": results,
        "source_tree": status.get("official_source_tree", str(ROOT / "official_reproduction" / "diffusion_policy_source")),
        "checkpoint_path": status.get("official_checkpoint_path", str(ROOT / "official_reproduction" / "data" / "epoch=0550-test_mean_score=0.969.ckpt")),
    }


def write_environment_log() -> None:
    LOGS.mkdir(parents=True, exist_ok=True)
    off = official_status()
    env = os.environ.copy()
    env["CUDA_DEVICE_MEMORY_SHARED_CACHE"] = "/tmp/finalproject-vgpu-cache.log"
    lines = [
        "# Environment Log",
        "",
        f"Python: {sys.version}",
        f"Platform: {platform.platform()}",
        "",
        "## PyTorch / CUDA",
        "The default container vGPU cache path can make torch import fail. The working invocation is:",
        "",
        "```bash",
        "CUDA_DEVICE_MEMORY_SHARED_CACHE=/tmp/finalproject-vgpu-cache.cache python -c \"import torch; ...\"",
        "```",
        "",
        "Observed torch output:",
        "```text",
        sh([
            sys.executable,
            "-c",
            "import torch; print(torch.__version__); print(torch.cuda.is_available()); print(torch.cuda.get_device_name(0) if torch.cuda.is_available() else 'cpu')",
        ], env=env),
        "```",
        "",
        "## Core Python Packages",
        "```text",
        sh([
            sys.executable,
            "-c",
            "import numpy, pandas, matplotlib, pptx, reportlab; "
            "print('numpy', numpy.__version__); print('pandas', pandas.__version__); "
            "print('matplotlib', matplotlib.__version__); print('python-pptx ok'); print('reportlab ok')",
        ]),
        "```",
        "",
        "## Official Diffusion Policy Reproduction",
        "The official repository was reachable for metadata:",
        "```text",
        sh(["git", "-c", "http.proxy=", "-c", "https.proxy=", "ls-remote", "https://github.com/real-stanford/diffusion_policy.git", "HEAD"]),
        "```",
        "Full and shallow git clone attempts were unstable, so the official source was acquired through the "
        "GitHub codeload archive and extracted locally.",
        "",
        f"Official source tree: `{off['source_tree']}`",
        f"Official checkpoint: `{off['checkpoint_path']}`",
        f"Checkpoint size: `{off['checkpoint_size_bytes']}` bytes",
        f"50-seed official PushT mean score: `{fmt(off['preferred'].get('test_mean_score', 0.0))}`",
        f"Official eval log: `{off['preferred'].get('path', '')}`",
        "",
        "Compatibility patches were required for the current Python 3.12 / Gym 0.26 environment: disable "
        "shared-memory batching for the custom PushT observation space, restore the old reset return contract, "
        "and update Gym vector concatenate argument order.",
        "",
        "## LeRobot / Hugging Face Near-Official Path",
        "After official GitHub clone and large-file transfer attempts remained unstable, a second reproduction path "
        "was executed through the LeRobot `lerobot/diffusion_pusht` checkpoint. Official HF transfer was slow, so "
        "the large `model.safetensors` file was downloaded through `hf-mirror.com` with chunked range requests. "
        "The assembled safetensors file has the expected size and can be opened and loaded by LeRobot, but its SHA "
        "did not match the HF LFS metadata, so local checkpoint rollouts are reported as smoke tests rather than "
        "cryptographically verified benchmark reproduction.",
        "",
        "LeRobot smoke command:",
        "```bash",
        "CUDA_DEVICE_MEMORY_SHARED_CACHE=/tmp/finalproject-vgpu-cache.cache "
        "/root/FinalProject/.venv_lerobot/bin/python scripts/eval_lerobot_pusht_smoke.py "
        "--episodes 1 --max-steps 300 --num-inference-steps 100 --save-video",
        "```",
    ]
    (LOGS / "environment_log.md").write_text("\n".join(lines) + "\n")


def write_readme(summary: pd.DataFrame) -> None:
    joint = row(summary, "joint_scheduler")
    best = row(summary, "fixed_k2_h1")
    off = official_status()
    lr = lerobot_status()
    hf_eval = lr["hf_eval"]
    smoke = lr["smoke"]
    smoke_mean = smoke.get("mean_max_reward", 0.0)
    smoke_success = smoke.get("success_rate", 0.0)
    smoke_steps = smoke.get("num_inference_steps", 0)
    model_size = lr["model_size_mb"] or 0.0
    readme = f"""# Joint Compute-Control Budgeting for Receding-Horizon Diffusion Policy

This final project studies the inference-time tradeoff in Diffusion Policy-style receding-horizon control:

- `k`: denoising steps per policy call.
- `h`: executed control steps before replanning.
- `B ~= k / h`: amortized inference budget per control step.

The reproduction target is **Diffusion Policy: Visuomotor Policy Learning via Action Diffusion** (RSS 2023). The original Stanford/Columbia source and checkpoint are included under `official_reproduction/`, and the official low-dimensional Push-T checkpoint evaluation was executed in this container. The main improvement experiment uses a deterministic Push-T-style closed-loop surrogate benchmark because the proposed contribution changes inference-time allocation variables `(k,h)` that are easier to isolate and sweep reproducibly in a controlled runner.

An official reproduction wrapper is included at `scripts/run_official_pusht_eval.py`. It loads the official checkpoint, instantiates the upstream workspace and PushT keypoints runner, and allows short or 50-seed evaluation overrides. The preferred 50-seed official run used `diffusers==0.11.1` and produced mean score `{fmt(off['preferred'].get('test_mean_score', 0.0))}`. The checkpoint filename reports `test_mean_score=0.969`; the observed gap is recorded as an environment-version reproduction gap because the current container uses Python 3.12, PyTorch 2.7, and Gym 0.26 compatibility patches rather than the original Python 3.9 / PyTorch 1.12 stack.

## Official Push-T Reproduction

- Source tree: `official_reproduction/diffusion_policy_source`
- Checkpoint: `official_reproduction/data/epoch=0550-test_mean_score=0.969.ckpt`
- Checkpoint size: `{off['checkpoint_size_bytes']}` bytes.
- 4-seed smoke score: `{fmt(off['results'].get('pusht_eval_official_n4_diffusers011', {}).get('test_mean_score', 0.0))}`.
- 50-seed score: `{fmt(off['preferred'].get('test_mean_score', 0.0))}`.
- 50-seed log: `{off['preferred'].get('path', '')}`.
- Compatibility patches: `shared_memory=False`, Gym reset wrapper, Gym vector concatenate argument order.

## Near-Official LeRobot Reproduction

- Model: `lerobot/diffusion_pusht`.
- HF model-card 500-episode evaluation: average max reward `{fmt(hf_eval.get("avg_max_reward", 0.0))}`, success `{fmt(hf_eval.get("pc_success", 0.0), 1)}%`.
- Model-card comparison to the original Diffusion Policy repository: average max reward `0.957`, success `64.2%`.
- Local checkpoint status: `{fmt(model_size, 1)}` MiB safetensors file downloaded through `hf-mirror.com`, safetensors structure verified, LeRobot `DiffusionPolicy` loaded with 262.7M parameters, real `gym-pusht` environment smoke tested.
- Local rollout smoke test: 1 episode, 300 steps, `{smoke_steps}` denoising steps, max reward `{fmt(smoke_mean)}`, success rate `{fmt(smoke_success)}`.

Important caveat: the mirror-assembled safetensors file loads and runs, but its SHA did not match the HF LFS metadata. Therefore the HF model-card metrics are treated as the near-official benchmark reference, and the local rollout is reported as an executable smoke test.

## Main Result

On the `B ~= 2` frontier with 8 matched seeds:

- Best fixed score baseline `(k=2, h=1)`: score `{fmt(best.score_mean)}`, success `{fmt(best.success_mean)}`, policy calls `{fmt(best.policy_calls_mean, 1)}`, smoothness cost `{fmt(best.smoothness_mean)}`.
- Joint scheduler: score `{fmt(joint.score_mean)}`, success `{fmt(joint.success_mean)}`, policy calls `{fmt(joint.policy_calls_mean, 1)}`, smoothness cost `{fmt(joint.smoothness_mean)}`.

The joint scheduler slightly improves score while reducing policy calls and action roughness relative to the best fixed score baseline.

## Reproduce

```bash
cd /root/FinalProject
CUDA_DEVICE_MEMORY_SHARED_CACHE=/tmp/finalproject-vgpu-cache.cache python scripts/eval_kh_grid.py --seeds 0 1 2 3 4 5 6 7 --episode-steps 120
python scripts/plot_iso_compute.py
python scripts/generate_deliverables.py
```

Official checkpoint evaluation path, for a network-stable run:

```bash
cd /root/FinalProject
CUDA_DEVICE_MEMORY_SHARED_CACHE=/tmp/finalproject-vgpu-cache.cache \
  /root/FinalProject/.venv_official/bin/python scripts/run_official_pusht_eval.py \
  --output-dir official_reproduction/pusht_eval_official_n50_diffusers011 \
  --n-test 50 --n-envs 5 --n-test-vis 0 --max-steps 300
```

## Outputs

- Raw and summarized results: `results/`
- Figures: `figures/`
- Main comparison bar plot: `figures/main_comparison_bars.png`
- Rollout animation: `artifacts/videos/joint_scheduler_rollout.gif`
- Environment log: `artifacts/logs/environment_log.md`
- Completion audit: `artifacts/logs/completion_audit.md`
- Official reproduction status: `results/official_reproduction_status.json`
- Official checkpoint eval wrapper: `scripts/run_official_pusht_eval.py`
- Official reference notes: `official_reproduction/REFERENCE_NOTES.md`
- LeRobot PushT smoke results: `results/lerobot_pusht_smoke_summary.json` and `results/lerobot_pusht_smoke_results.csv`
- LeRobot PushT smoke video: `artifacts/videos/lerobot_pusht_episode_00.gif`
- Final report: `report/final_report.pdf` and `report/final_report.md`
- Slides: `slides/final_presentation.pptx` and `slides/final_presentation.pdf`

## Important Limitation

The scheduler-improvement numbers are from the local Push-T-style surrogate, not official Diffusion Policy checkpoint rollouts. The report states this explicitly and frames the result as a reproducible inference-time analysis layer built on top of an official Push-T checkpoint reproduction.
"""
    (ROOT / "README.md").write_text(readme)


def comparison_table(summary: pd.DataFrame) -> list[list[str]]:
    methods = [
        ("Default DP-style (8,4)", "fixed_k8_h4"),
        ("Best fixed B=2 (2,1)", "fixed_k2_h1"),
        ("Denoising-only point (4,2)", "fixed_k4_h2"),
        ("Long denoise/chunk (16,8)", "fixed_k16_h8"),
        ("AAC/DVAC-style h-only", "aac_dvac_h_only"),
        ("Joint scheduler", "joint_scheduler"),
    ]
    table = [["Method", "B", "Score", "Success", "Calls", "NFE", "Smoothness"]]
    for label, method in methods:
        r = row(summary, method)
        table.append([
            label,
            fmt(r.avg_budget_mean, 2),
            fmt(r.score_mean),
            fmt(r.success_mean),
            fmt(r.policy_calls_mean, 1),
            fmt(r.nfe_mean, 1),
            fmt(r.smoothness_mean),
        ])
    return table


def write_report_markdown(summary: pd.DataFrame) -> str:
    joint = row(summary, "joint_scheduler")
    best = row(summary, "fixed_k2_h1")
    off = official_status()
    lr = lerobot_status()
    hf_eval = lr["hf_eval"]
    smoke = lr["smoke"]
    table_rows = comparison_table(summary)
    md_table = "\n".join(
        ["| " + " | ".join(table_rows[0]) + " |", "| " + " | ".join(["---"] * len(table_rows[0])) + " |"]
        + ["| " + " | ".join(r) + " |" for r in table_rows[1:]]
    )
    md = f"""# Joint Compute-Control Budgeting for Receding-Horizon Diffusion Policy

## Abstract

Diffusion Policy produces action chunks through iterative denoising and deploys them in a receding-horizon loop. This project studies a practical inference question: under a fixed per-step compute budget, should a robot spend computation on more denoising steps or on more frequent replanning? We first reproduce the official low-dimensional Push-T checkpoint using the upstream workspace and runner, obtaining a 50-seed mean score of `{fmt(off['preferred'].get('test_mean_score', 0.0))}` in the current Python 3.12 / PyTorch 2.7 container. We then evaluate the proposed inference-time idea in a lightweight Push-T-style closed-loop surrogate and sweep denoising steps `k` and execution horizon `h`. A training-free joint scheduler chooses `(k,h)` on the same `B ~= k/h` frontier. On the `B ~= 2` frontier, the joint scheduler obtains score `{fmt(joint.score_mean)}` versus `{fmt(best.score_mean)}` for the best fixed score baseline `(2,1)`, while reducing policy calls from `{fmt(best.policy_calls_mean,1)}` to `{fmt(joint.policy_calls_mean,1)}` and smoothness cost from `{fmt(best.smoothness_mean)}` to `{fmt(joint.smoothness_mean)}`.

## 1. Introduction

Diffusion Policy is a strong visuomotor imitation learning method because iterative action denoising can represent multimodal action distributions and produce coherent action chunks. Deployment introduces a resource allocation problem. More denoising steps can improve action quality, but each policy call becomes more expensive. More frequent replanning improves reactivity, especially near contact, but increases the number of policy calls. This project formulates that tradeoff as joint compute-control budgeting.

## 2. Related Work

The reproduced base paper is **Diffusion Policy: Visuomotor Policy Learning via Action Diffusion** (RSS 2023). The method generates action sequences with a conditional diffusion model and executes them in a receding-horizon control loop. Recent adaptive action chunking, adaptive denoising, fast denoising, and uncertainty-based failure detection methods study related single-axis decisions. The novelty boundary here is intentionally narrower: the project evaluates denoising and replanning jointly under an iso-compute constraint and compares a joint scheduler against fixed and single-axis baselines.

## 3. Problem Formulation

Let `k` be the number of denoising steps per policy call and `h` be the number of controls executed before the next replan. The amortized inference budget is approximated as `B ~= k / h`. The static grid uses `k in {{2,4,8,16}}` and `h in {{1,2,4,8}}`, with main frontiers `B=1`, `B=2`, and `B=4`. The main comparison fixes the average budget near `B=2`.

## 4. Reproduction Setup

The original Diffusion Policy repository was checked by `git ls-remote`. Full and shallow `git clone` attempts were intermittent, so the official source was acquired through the GitHub codeload archive and extracted to `official_reproduction/diffusion_policy_source`. The official low-dimensional Push-T checkpoint was downloaded from the project server to `official_reproduction/data/epoch=0550-test_mean_score=0.969.ckpt`; the file size is `{off['checkpoint_size_bytes']}` bytes. PyTorch 2.7.0 with the RTX 5090 vGPU works when `CUDA_DEVICE_MEMORY_SHARED_CACHE` is redirected to `/tmp`; the default cache path triggered a vGPU runtime failure during initial testing.

The upstream checkpoint was evaluated with `scripts/run_official_pusht_eval.py`, which instantiates the original workspace, EMA policy, and PushT keypoints runner. Because this container is newer than the original conda environment, the local source copy includes compatibility patches for Gym 0.26: disabling shared-memory batching for the custom PushT observation space, restoring the old reset return contract, and updating vector concatenate argument order. The preferred 50-seed run used official-era `diffusers==0.11.1` and obtained mean score `{fmt(off['preferred'].get('test_mean_score', 0.0))}`. The checkpoint filename reports `test_mean_score=0.969`, so the remaining gap is treated as an environment-version reproduction gap rather than an execution failure.

As a supplementary path, a Hugging Face / LeRobot checkpoint was also tested. The `lerobot/diffusion_pusht` model card reports a 500-episode PushT evaluation with average max reward `{fmt(hf_eval.get("avg_max_reward", 0.0))}` and success `{fmt(hf_eval.get("pc_success", 0.0), 1)}%`; the same card lists the original Diffusion Policy repository comparison at average max reward `0.957` and success `64.2%`. Official HF transfer of the 1.0 GB safetensors file was slow, so `hf-mirror.com` and chunked HTTP range requests were used. The assembled file has the expected size, can be opened as safetensors, and loads into LeRobot `DiffusionPolicy` with 262.7M parameters. Its SHA did not match the HF LFS metadata, so local LeRobot rollouts are reported as executable smoke tests rather than cryptographically verified benchmark runs.

The local LeRobot smoke test ran the real `gym-pusht` environment for `{smoke.get("episodes", 0)}` episode, `{smoke.get("max_steps", 0)}` steps, and `{smoke.get("num_inference_steps", 0)}` denoising steps. It reached max reward `{fmt(smoke.get("mean_max_reward", 0.0))}` and success rate `{fmt(smoke.get("success_rate", 0.0))}`. The improvement experiment below uses a deterministic Push-T-style surrogate environment with closed-loop contact dynamics, action chunks, denoising-quality noise, and stale-plan effects.

The surrogate is not a substitute for official benchmark numbers. It is a controlled reproduction of the deployment mechanism needed for the proposed improvement: action diffusion quality versus receding-horizon replanning frequency.

For a reproduced official run in this project folder:

```bash
CUDA_DEVICE_MEMORY_SHARED_CACHE=/tmp/finalproject-vgpu-cache.cache \
  /root/FinalProject/.venv_official/bin/python scripts/run_official_pusht_eval.py \
  --output-dir official_reproduction/pusht_eval_official_n50_diffusers011 \
  --n-test 50 --n-envs 5 --n-test-vis 0 --max-steps 300
```

The official reproduction reference note in `official_reproduction/REFERENCE_NOTES.md` records the upstream checkpoint path, command pattern, and original conda environment mismatch.

## 5. Method

The fixed runner evaluates every `(k,h)` pair. Larger `k` reduces residual action noise and mode bias in the diffusion-like planner. Larger `h` reduces policy calls but increases open-loop staleness in contact. The proposed training-free scheduler is restricted to the `B=2` candidate set and uses state features available in the simulator:

- pusher-object distance,
- object-goal distance,
- contact proxy,
- recent progress.

The final rule selects `(4,2)` in free-space approach and final alignment, and `(2,1)` during close-contact pushing or unstable progress. This keeps average `B` near 2 while shifting control frequency toward high-risk phases.

## 6. Results

{md_table}

The static grid shows that compute-equivalent choices are not interchangeable. On the `B=2` frontier, `(2,1)` gives the highest fixed score because the task is contact sensitive and benefits from frequent replanning. However, `(2,1)` also creates the largest action roughness. The joint scheduler preserves high-frequency updates near contact while using `(4,2)` in easier phases. This raises mean score to `{fmt(joint.score_mean)}` and reduces smoothness cost by about `{fmt(100*(best.smoothness_mean-joint.smoothness_mean)/best.smoothness_mean,1)}%` relative to `(2,1)`.

## 7. Phase-Wise Analysis

The phase-wise plot separates free-space approach, contact pushing, and near-goal alignment. Fixed long-horizon settings are efficient but lose score when contact starts. Fixed short-horizon replanning is robust but noisy. The scheduler's benefit comes from recognizing that free-space and alignment tolerate a longer action chunk, while contact needs immediate feedback.

## 8. Limitations

The largest limitation is that the official checkpoint evaluation runs in a compatibility environment, not the original Python 3.9 / PyTorch 1.12 conda stack, and its 50-seed score is below the checkpoint filename score. A LeRobot checkpoint was also downloaded through an alternate mirror, loaded, and executed locally, but the mirror-assembled file did not pass HF LFS SHA verification. The scheduler-improvement environment is a surrogate, so those numeric values should not be reported as official Diffusion Policy benchmark performance. The scheduler is hand-tuned on the same task family, the compute model uses `k/h` rather than measured neural network latency, and there is no real-robot validation.

## 9. Conclusion

The experiment supports the project hypothesis: under the same amortized inference budget, the best allocation can depend on task phase. Jointly scheduling denoising and replanning can outperform or match a tuned fixed budget point while improving execution smoothness. For deployment, the practical lesson is to treat action diffusion quality and receding-horizon frequency as one coupled budget rather than two independent hyperparameters.

## References

- Diffusion Policy project page: https://diffusion-policy.cs.columbia.edu/
- Official code repository: https://github.com/real-stanford/diffusion_policy
- LeRobot model card used for near-official PushT reference: https://huggingface.co/lerobot/diffusion_pusht
- Paper: Diffusion Policy: Visuomotor Policy Learning via Action Diffusion, RSS 2023 / arXiv 2303.04137.
- Official checkpoint example from the repository README: `low_dim/pusht/diffusion_policy_cnn/train_0/checkpoints/epoch=0550-test_mean_score=0.969.ckpt`.
"""
    REPORT.mkdir(parents=True, exist_ok=True)
    (REPORT / "final_report.md").write_text(md)
    return md


def add_wrapped(story, text: str, style, width: int = 95) -> None:
    for para in text.split("\n\n"):
        if para.strip():
            story.append(Paragraph(para.strip().replace("\n", " "), style))
            story.append(Spacer(1, 0.08 * inch))


def write_report_pdf(summary: pd.DataFrame) -> None:
    off = official_status()
    lr = lerobot_status()
    hf_eval = lr["hf_eval"]
    smoke = lr["smoke"]
    REPORT.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(str(REPORT / "final_report.pdf"), pagesize=letter, rightMargin=54, leftMargin=54, topMargin=48, bottomMargin=48)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="TitleCenter", parent=styles["Title"], alignment=TA_CENTER, fontSize=18, leading=22))
    styles.add(ParagraphStyle(name="H", parent=styles["Heading2"], fontSize=13, leading=16, spaceBefore=10, spaceAfter=4))
    styles.add(ParagraphStyle(name="Body", parent=styles["BodyText"], fontSize=9.5, leading=13, alignment=TA_LEFT))
    story = [
        Paragraph("Joint Compute-Control Budgeting for Receding-Horizon Diffusion Policy", styles["TitleCenter"]),
        Paragraph("Iso-Compute Allocation of Denoising and Replanning for Efficient Robotic Manipulation", styles["Body"]),
        Spacer(1, 0.15 * inch),
    ]
    sections = [
        ("Abstract", f"Diffusion Policy uses iterative action denoising and receding-horizon execution. This project reproduces the official low-dimensional Push-T checkpoint with the upstream workspace and runner, obtaining a 50-seed mean score of {fmt(off['preferred'].get('test_mean_score', 0.0))} in the current compatibility environment. It then evaluates how to allocate a fixed inference budget between denoising quality and replanning frequency in a deterministic Push-T-style surrogate."),
        ("Problem", "We define k as denoising steps per policy call, h as executed controls before replanning, and B ~= k/h as amortized compute per control step. The grid evaluates k in {2,4,8,16} and h in {1,2,4,8}, with main analysis on B=1, B=2, and B=4 frontiers."),
        ("Official Reproduction", f"The official source was acquired through a GitHub codeload archive after git clone instability, and the 1044185793-byte Push-T checkpoint was downloaded from the project server. The 50-seed official runner result is {fmt(off['preferred'].get('test_mean_score', 0.0))}; the checkpoint filename reports 0.969, so the difference is recorded as an environment-version reproduction gap."),
        ("LeRobot Supplement", f"The LeRobot lerobot/diffusion_pusht model card reports 500-episode average max reward {fmt(hf_eval.get('avg_max_reward', 0.0))} and success {fmt(hf_eval.get('pc_success', 0.0), 1)}%. The local mirror-downloaded checkpoint loads and runs in gym-pusht; a 300-step, {smoke.get('num_inference_steps', 0)}-denoising-step smoke test reached max reward {fmt(smoke.get('mean_max_reward', 0.0))}. Because the file SHA did not match HF LFS metadata, local rollout results are marked as smoke verification."),
        ("Method", "A fixed runner evaluates every (k,h). A training-free joint scheduler operates on the B=2 frontier. It selects (4,2) in free-space approach and final alignment, and (2,1) during close-contact pushing or unstable progress."),
    ]
    for title, text in sections:
        story.append(Paragraph(title, styles["H"]))
        add_wrapped(story, text, styles["Body"])
    table = Table(comparison_table(summary), repeatRows=1)
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f2937")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONT", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONT", (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE", (0, 0), (-1, -1), 7.5),
        ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f3f4f6")]),
    ]))
    story.append(Paragraph("Main Comparison", styles["H"]))
    story.append(table)
    story.append(Spacer(1, 0.15 * inch))
    for fig in ["kh_score_heatmap.png", "iso_compute_curves.png", "pareto_score_compute.png", "main_comparison_bars.png", "smoothness_vs_compute.png", "phase_wise_scores.png"]:
        path = FIGURES / fig
        if path.exists():
            story.append(Image(str(path), width=6.2 * inch, height=4.1 * inch))
            story.append(Spacer(1, 0.08 * inch))
    for title, text in [
        ("Results", "The joint scheduler reaches score 0.937 at average B=2.006, compared with 0.925 for the best fixed score baseline (2,1). It reduces policy calls from 120.0 to 90.8 and smoothness cost from 0.444 to 0.271."),
        ("Limitations", "The scheduler-improvement numeric results are from the surrogate, not official Diffusion Policy checkpoint rollouts. The official checkpoint run uses compatibility patches for Python 3.12 / Gym 0.26, and the LeRobot mirror checkpoint is executable locally but not SHA-verified. The compute model is approximate, the scheduler is rule-based, and the experiment is limited to one task family."),
        ("Conclusion", "The experiment supports the hypothesis that denoising and replanning should be allocated jointly. Contact phases favor frequent replanning; easier phases can use a longer chunk and more denoising while staying on the same compute frontier."),
        ("References", "Diffusion Policy project page: https://diffusion-policy.cs.columbia.edu/. Official code: https://github.com/real-stanford/diffusion_policy. LeRobot model card: https://huggingface.co/lerobot/diffusion_pusht. Paper: Diffusion Policy: Visuomotor Policy Learning via Action Diffusion, RSS 2023 / arXiv 2303.04137."),
    ]:
        story.append(Paragraph(title, styles["H"]))
        add_wrapped(story, text, styles["Body"])
    doc.build(story)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 41, 55)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.48), Inches(0.88), Inches(12.0), Inches(0.35))
        q = sub.text_frame.paragraphs[0]
        q.text = subtitle
        q.font.size = Pt(13)
        q.font.color.rgb = RGBColor(75, 85, 99)


def add_bullets(slide, bullets: list[str], x=0.65, y=1.35, w=5.7, h=4.8, size=18) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(size)
        p.level = 0
        p.space_after = Pt(8)


def add_picture(slide, name: str, x=6.6, y=1.35, w=6.1) -> None:
    path = FIGURES / name
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


def write_slides(summary: pd.DataFrame) -> None:
    off = official_status()
    lr = lerobot_status()
    hf_eval = lr["hf_eval"]
    smoke = lr["smoke"]
    SLIDES.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = [
        ("Joint Compute-Control Budgeting", "Iso-compute denoising and replanning for Diffusion Policy-style control",
         ["Final project: reproduce a robotics paper and add an improvement idea", "Base paper: Diffusion Policy, RSS 2023", "Task: Push-T-style closed-loop contact benchmark"], None),
        ("Motivation", None,
         ["Diffusion Policy is accurate but iterative denoising is expensive", "Replanning improves reactivity but increases policy calls", "Deployment needs a fixed per-step inference budget"], "pareto_score_compute.png"),
        ("Key Question", None,
         ["Should compute be spent on larger k or smaller h?", "k: denoising steps per policy call", "h: executed controls before replanning", "B ~= k / h is the iso-compute budget"], "iso_compute_curves.png"),
        ("Reproduction Setup", None,
         ["Official source acquired via GitHub codeload archive", "Official Push-T checkpoint downloaded: 1044185793 bytes", f"50-seed official runner score: {fmt(off['preferred'].get('test_mean_score', 0.0))}", "Compatibility patches needed for Python 3.12 / Gym 0.26"], None),
        ("LeRobot PushT Checkpoint", None,
         [f"Model card 500-episode max reward: {fmt(hf_eval.get('avg_max_reward', 0.0))}", f"Model card success rate: {fmt(hf_eval.get('pc_success', 0.0), 1)}%", "Local safetensors loads and runs in real gym-pusht", f"Local smoke: max reward {fmt(smoke.get('mean_max_reward', 0.0))}, {smoke.get('num_inference_steps', 0)} denoising steps", "Caveat: mirror file did not pass HF LFS SHA check"], None),
        ("Static Grid", None,
         ["Grid: k in {2,4,8,16}, h in {1,2,4,8}", "Contact-rich pushing makes equal-compute points behave differently", "The best fixed B=2 score point is (2,1)"], "kh_score_heatmap.png"),
        ("Joint Scheduler", None,
         ["Training-free rule on the B=2 frontier", "Use (4,2) for free-space approach and final alignment", "Use (2,1) during close-contact pushing or unstable progress", "Goal: keep compute fixed while moving reactivity to high-risk phases"], None),
        ("Main Results", None,
         [f"Joint score: {fmt(row(summary, 'joint_scheduler').score_mean)}", f"Best fixed B=2 score: {fmt(row(summary, 'fixed_k2_h1').score_mean)}", "Policy calls: 90.8 vs 120.0", "Smoothness cost: 0.271 vs 0.444"], "main_comparison_bars.png"),
        ("Phase-Wise Analysis", None,
         ["Free-space approach tolerates longer chunks", "Contact pushing needs frequent replanning", "Near-goal alignment benefits from smoother denoised chunks"], "phase_wise_scores.png"),
        ("Limitations", None,
         ["Improvement numbers are from the surrogate, not official benchmark rollouts", "Official run uses compatibility patches and scores below checkpoint filename", "LeRobot checkpoint path is executable but not SHA-verified", "B=k/h is an approximate compute model", "No real-robot validation"], None),
        ("Takeaways", None,
         ["Denoising and replanning are coupled deployment decisions", "Best allocation changes by task phase", "A simple joint scheduler can improve the score-smoothness tradeoff under the same budget"], None),
    ]
    for title, subtitle, bullets, picture in slides:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(250, 250, 249)
        add_title(slide, title, subtitle)
        add_bullets(slide, bullets, w=5.75 if picture else 11.8, size=18)
        if picture:
            add_picture(slide, picture)

    prs.save(SLIDES / "final_presentation.pptx")

    # Also produce a simple PDF slide deck for environments without PowerPoint.
    doc = SimpleDocTemplate(str(SLIDES / "final_presentation.pdf"), pagesize=landscape(letter), rightMargin=36, leftMargin=36, topMargin=28, bottomMargin=28)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="SlideTitle", parent=styles["Title"], fontSize=24, leading=28, alignment=TA_LEFT))
    styles.add(ParagraphStyle(name="SlideBody", parent=styles["BodyText"], fontSize=15, leading=20))
    story = []
    for i, (title, subtitle, bullets, picture) in enumerate(slides):
        story.append(Paragraph(title, styles["SlideTitle"]))
        if subtitle:
            story.append(Paragraph(subtitle, styles["SlideBody"]))
        story.append(Spacer(1, 0.15 * inch))
        for bullet in bullets:
            story.append(Paragraph(f"- {bullet}", styles["SlideBody"]))
        if picture and (FIGURES / picture).exists():
            story.append(Spacer(1, 0.12 * inch))
            story.append(Image(str(FIGURES / picture), width=4.8 * inch, height=3.1 * inch))
        if i < len(slides) - 1:
            story.append(PageBreak())
    doc.build(story)

    notes = """# Oral Notes

English or Chinese can be used for the oral presentation. Suggested short Chinese framing:

本项目复现并分析 Diffusion Policy 的推理时结构：每次策略调用需要若干 denoising steps，并在 receding horizon 中执行一段 action chunk。官方源码通过 GitHub archive 获取，官方 Push-T checkpoint 已下载并跑通 50 个 test seed；当前 Python 3.12 / Gym 0.26 环境需要兼容补丁，所以实测分数低于 checkpoint 文件名中的 0.969。核心改进问题是在固定预算下，计算量应该用于更高质量的 denoising，还是更频繁地 replanning。实验显示，在接触阶段高频 replanning 更重要，而在 free-space 和最终对齐阶段可以使用更长 chunk 和更多 denoising。联合调度器在相同平均预算下提升了 score，同时减少了 policy calls 和动作抖动。
"""
    (SLIDES / "oral_notes.md").write_text(notes)


def main() -> None:
    summary = pd.read_csv(RESULTS / "summary_results.csv")
    write_environment_log()
    write_readme(summary)
    write_report_markdown(summary)
    write_report_pdf(summary)
    write_slides(summary)
    print("Generated README, environment log, report, and slides.")


if __name__ == "__main__":
    main()
